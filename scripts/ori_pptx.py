#!/usr/bin/env python3
"""Ori-native PPTX generator (see references/office.md).

Mirrors templates/slides.html: evidence/statement layouts, meta rail with
the mark, exact thread rule (node + blue segment + hairline), pinned
takeaway, night cover/close. Sizes derive from the 1280×720 template at
96dpi (1px = 0.75pt; 13.333×7.5in canvas).
Running this file directly writes ori-pptx-demo.pptx as a visual check.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Ori tokens ───────────────────────────────────────────────────────
BLUE = RGBColor(0x2F, 0x80, 0xED)
NAVY = RGBColor(0x21, 0x3C, 0x60)
INK = RGBColor(0x00, 0x04, 0x16)
SLATE = RGBColor(0x4A, 0x4D, 0x5A)
MIST = RGBColor(0x8A, 0x94, 0xAB)
FROST = RGBColor(0xFA, 0xFB, 0xFE)
PAPER = RGBColor(0xFF, 0xFF, 0xFF)
WASH = RGBColor(0xEF, 0xF5, 0xFD)
LINE = RGBColor(0xE3, 0xE8, 0xF4)
LINE_STRONG = RGBColor(0xC9, 0xD2, 0xE8)
NIGHT = RGBColor(0x02, 0x12, 0x2A)
LINE_N = RGBColor(0x14, 0x29, 0x4D)
INK_N = RGBColor(0xF4, 0xF7, 0xFF)
MIST_N = RGBColor(0x8F, 0xA0, 0xC2)
BLUE_SOFT = RGBColor(0x7E, 0xAF, 0xF1)
BLUE_FAINT = RGBColor(0xB6, 0xD1, 0xF7)
TINT = RGBColor(0xE3, 0xEE, 0xFC)
CHART_RAMP = [BLUE, NAVY, BLUE_SOFT, BLUE_FAINT, TINT]

SANS = "Inter"
MONO = "IBM Plex Mono"

ASSETS = Path(__file__).resolve().parent.parent / "assets" / "logo" / "png"

W, H = Inches(13.333), Inches(7.5)
MX = Inches(0.75)            # 72px side padding
CW = Inches(13.333 - 1.5)    # content width


def px(v):                   # template px → EMU (1280px canvas = 13.333in)
    return Emu(int(v * 914400 / 96))


class OriDeck:
    def __init__(self, occasion="Deck", date="", classification=""):
        self.prs = Presentation()
        self.prs.slide_width, self.prs.slide_height = W, H
        self.blank = self.prs.slide_layouts[6]
        self.occasion, self.date = occasion, date
        self.classification = classification
        self.n = 0
        self.total = "NN"      # patched on save

    # ── primitives ──
    def _slide(self, night=False):
        s = self.prs.slides.add_slide(self.blank)
        s.background.fill.solid()
        s.background.fill.fore_color.rgb = NIGHT if night else FROST
        self.n += 1
        return s

    def _rect(self, s, x, y, w, h, fill):
        from pptx.enum.shapes import MSO_SHAPE
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        r.fill.solid()
        r.fill.fore_color.rgb = fill
        r.line.fill.background()
        r.shadow.inherit = False
        return r

    def _text(self, s, x, y, w, h, runs, align=PP_ALIGN.LEFT,
              line_spacing=None, space_after=0):
        """runs: [(text, font, size_pt, color, bold, caps)] per paragraph,
        or a list of such lists for multiple paragraphs."""
        tb = s.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        paragraphs = runs if isinstance(runs[0], list) else [runs]
        for i, para in enumerate(paragraphs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if line_spacing:
                p.line_spacing = line_spacing
            p.space_after = Pt(space_after)
            for text, font, size, color, bold, caps in para:
                r = p.add_run()
                r.text = text.upper() if caps else text
                r.font.name, r.font.size = font, Pt(size)
                r.font.color.rgb, r.font.bold = color, bold
        return tb

    def _rail(self, s, night=False, pageno=True):
        c = MIST_N if night else MIST
        brand_c = INK_N if night else INK
        mark = ASSETS / ("mark-white.png" if night else "mark.png")
        y = px(56)
        if mark.exists():
            s.shapes.add_picture(str(mark), MX, y, height=px(16))
        left = [("SOCIALTRAIT", MONO, 9, brand_c, True, True),
                (f"  ·  {self.occasion}", MONO, 9, c, False, True)]
        if self.date and self.n == 1:
            left.append((f"  ·  {self.date}", MONO, 9, c, False, True))
        self._text(s, MX + px(26), y, px(700), px(20), left)
        right = (self.classification.upper() if self.n == 1
                 else f"{self.n:02d} / {self.total}")
        if pageno:
            self._text(s, W - MX - px(300), y, px(300), px(20),
                       [(right, MONO, 9, c, False, True)],
                       align=PP_ALIGN.RIGHT)
        self._rect(s, MX, y + px(28), CW, px(1), LINE_N if night else LINE)

    def _thread(self, s, y):
        self._rect(s, MX, y, CW, px(1.5), LINE)          # hairline
        self._rect(s, MX, y, px(48), px(2), BLUE)        # blue segment
        self._rect(s, MX, y - px(2.5), px(7), px(7), BLUE)  # node

    def _eyebrow(self, s, y, idx, text, night=False):
        runs = [(f"{idx:02d}  ", MONO, 10, MIST_N if night else MIST,
                 False, True),
                (text, MONO, 10, BLUE_SOFT if night else BLUE, False, True)]
        self._text(s, MX, y, px(700), px(22), runs)

    def takeaway(self, s, text, label="Takeaway"):
        y, h = H - px(48) - px(62), px(62)
        self._rect(s, MX, y, CW, h, WASH)
        self._rect(s, MX, y, px(3), h, BLUE)
        self._text(s, MX + px(18), y + px(10),
                   CW - px(36), h - px(20),
                   [[(label, MONO, 8, BLUE, False, True)],
                    [(text, SANS, 13, INK, False, False)]])

    # ── slide archetypes ──
    def cover(self, eyebrow, title, lede, presenter=""):
        s = self._slide(night=True)
        self._rail(s, night=True, pageno=bool(self.classification))
        self._eyebrow(s, px(200), 0, eyebrow, night=True)
        # cover eyebrow has no index — rewrite without it
        s.shapes[-1].text_frame.paragraphs[0].runs[0].text = ""
        self._text(s, MX, px(232), px(760), px(220),
                   [(title, SANS, 40, INK_N, True, False)], line_spacing=1.05)
        self._text(s, MX, px(470), px(640), px(80),
                   [(lede, SANS, 16, MIST_N, False, False)], line_spacing=1.4)
        word = ASSETS / "wordmark-white.png"
        if word.exists():
            s.shapes.add_picture(str(word), MX, H - px(56) - px(24),
                                 height=px(22))
        if presenter:
            self._text(s, W - MX - px(400), H - px(56) - px(22), px(400),
                       px(20), [(presenter, MONO, 9, MIST_N, False, True)],
                       align=PP_ALIGN.RIGHT)
        return s

    def evidence(self, idx, section, title, context_stat=None,
                 context_kicker=None):
        """context_stat: (label, value, unit) · context_kicker: str"""
        s = self._slide()
        self._rail(s)
        self._eyebrow(s, px(104), idx, section)
        self._text(s, MX, px(132), px(660), px(80),
                   [(title, SANS, 22, INK, True, False)], line_spacing=1.15)
        cx = MX + px(710)
        if context_stat:
            label, value, unit = context_stat
            self._text(s, cx, px(136), px(420), px(20),
                       [(label, MONO, 8, MIST, False, True)])
            self._text(s, cx, px(160), px(420), px(46),
                       [(value, MONO, 30, INK, True, False),
                        (f" {unit}", MONO, 15, MIST, False, False)])
        elif context_kicker:
            self._text(s, cx, px(140), px(420), px(66),
                       [(context_kicker, SANS, 11, SLATE, False, False)],
                       line_spacing=1.4)
        self._thread(s, px(224))
        return s

    def statement(self, idx, section, title, lede=""):
        s = self._slide()
        self._rail(s)
        self._eyebrow(s, px(180), idx, section)
        self._text(s, MX, px(212), px(860), px(160),
                   [(title, SANS, 33, INK, True, False)], line_spacing=1.1)
        if lede:
            self._text(s, MX, px(390), px(660), px(90),
                       [(lede, SANS, 15, SLATE, False, False)],
                       line_spacing=1.45)
        return s

    def close(self, eyebrow, title, lede, contact=""):
        s = self._slide(night=True)
        self._rail(s, night=True)
        self._eyebrow(s, px(220), 0, eyebrow, night=True)
        s.shapes[-1].text_frame.paragraphs[0].runs[0].text = ""
        self._text(s, MX, px(252), px(860), px(160),
                   [(title, SANS, 36, INK_N, True, False)], line_spacing=1.08)
        self._text(s, MX, px(430), px(640), px(60),
                   [(lede, SANS, 15, MIST_N, False, False)], line_spacing=1.4)
        if contact:
            self._text(s, MX, H - px(56) - px(22), px(600), px(20),
                       [(contact, MONO, 9, MIST_N, False, True)])
        return s

    # ── evidence bodies ──
    BODY_Y = 264   # px — content starts under the thread

    def metric_row(self, s, metrics, y=None):
        y = px(y if y is not None else self.BODY_Y + 40)
        n = len(metrics)
        gap = px(1)
        cw = int((CW - gap * (n - 1)) / n)
        for i, m in enumerate(metrics):
            x = MX + i * (cw + gap)
            card = self._rect(s, x, y, cw, px(170), PAPER)
            card.line.color.rgb = LINE
            card.line.width = Pt(0.75)
            self._text(s, x + px(24), y + px(22), cw - px(48), px(18),
                       [(m["label"], MONO, 9, MIST, False, True)])
            self._text(s, x + px(24), y + px(52), cw - px(48), px(56),
                       [(m["value"], MONO, 38, BLUE if m.get("hero") else INK,
                         True, False),
                        (f" {m.get('unit','')}", MONO, 18, MIST, False,
                         False)])
            self._text(s, x + px(24), y + px(120), cw - px(48), px(36),
                       [(m.get("note", ""), SANS, 11, SLATE, False, False)])

    def bullets(self, s, items, x=None, y=None, w=None, size=14):
        x = x if x is not None else MX
        tb = s.shapes.add_textbox(x, px(y or self.BODY_Y + 30),
                                  w or px(1000), px(300))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.line_spacing, p.space_after = 1.35, Pt(10)
            r = p.add_run()
            r.text = "▪  "
            r.font.name, r.font.size, r.font.color.rgb = SANS, Pt(size), BLUE
            r2 = p.add_run()
            r2.text = item
            r2.font.name, r2.font.size = SANS, Pt(size)
            r2.font.color.rgb = SLATE

    def bar_chart(self, s, categories, series, x=None, y=None,
                  w=None, h=None):
        """series: [(name, [values])] — native, editable in Slides."""
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
        data = CategoryChartData()
        data.categories = categories
        for name, vals in series:
            data.add_series(name, vals)
        gf = s.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED,
            x if x is not None else MX, px(y or self.BODY_Y + 20),
            w or px(620), h or px(330), data)
        chart = gf.chart
        chart.has_title = False
        for i, srs in enumerate(chart.series):
            srs.format.fill.solid()
            srs.format.fill.fore_color.rgb = CHART_RAMP[i % len(CHART_RAMP)]
        chart.has_legend = len(series) > 1
        if chart.has_legend:
            chart.legend.position = XL_LEGEND_POSITION.TOP
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
            chart.legend.font.name = MONO
        for ax in (chart.category_axis, chart.value_axis):
            ax.tick_labels.font.size = Pt(10)
            ax.tick_labels.font.name = MONO
            ax.tick_labels.font.color.rgb = MIST
            ax.format.line.color.rgb = LINE_STRONG
        chart.value_axis.has_major_gridlines = True
        chart.value_axis.major_gridlines.format.line.color.rgb = LINE
        return gf

    def save(self, path):
        # patch page totals in rails
        total = f"{self.n:02d}"
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame and "/ NN" in shape.text_frame.text:
                    for p in shape.text_frame.paragraphs:
                        for r in p.runs:
                            r.text = r.text.replace("NN", total)
        self.prs.save(path)
        return path


# ── demo / visual regression check ──────────────────────────────────
if __name__ == "__main__":
    d = OriDeck(occasion="Q3 review", date="2026-07-08",
                classification="Confidential")
    d.cover("Pilot results", "Simulation predicted the market. "
            "Three times out of three.",
            "Q3 audience-simulation pilot results and the case for a "
            "ten-account expansion.", "Suraj · Socialtrait")
    s = d.evidence(1, "Headline numbers",
                   "Three pilots, one pattern: the simulation called "
                   "every winner",
                   context_kicker="Enterprise pilots across CPG, fintech, "
                   "and retail — 12 audience cells, Apr–Jun 2026.")
    d.metric_row(s, [
        {"label": "Prediction accuracy", "value": "87", "unit": "%",
         "note": "interval coverage 11/12 cells", "hero": True},
        {"label": "Time to insight", "value": "4", "unit": "hrs",
         "note": "panels took 19 days median"},
        {"label": "Media saved", "value": "$340", "unit": "k",
         "note": "one flagged creative, one pilot"},
    ])
    d.takeaway(s, "Simulation is accurate enough to gate media spend today.")
    s = d.evidence(2, "Speed", "Answers land inside the campaign cycle, "
                   "not after it",
                   context_stat=("Median time to insight", "4", "hrs"))
    d.bar_chart(s, ["Field study", "Online panel", "Socialtrait"],
                [("days to insight", [28, 19, 0.17])], w=px(620), h=px(300))
    d.bullets(s, ["9-day median creative iteration cycle at pilot clients",
                  "Panels answer after the decision; simulation answers "
                  "before",
                  "Same-day re-tests after every creative revision"],
              x=MX + px(700), w=px(430), size=12)
    d.takeaway(s, "Research finally moves at the speed the campaign "
               "already does.")
    d.close("Decision requested", "Expand the pilot. Ten accounts, "
            "this quarter.", "Approval needed by July 18 · detail in "
            "ST-OP-014.", "suraj@socialtrait.ai")
    print("wrote", d.save(str(Path.cwd() / "ori-pptx-demo.pptx")))
